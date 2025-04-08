import os
import streamlit as st
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from openai import OpenAI

# Load environment variables
load_dotenv()

# Set up OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Extract text from PPTX slides
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text

# Generate summary and extracted sections using OpenAI GPT
def summarize_and_extract(text):
    prompt = f"""
    From the following incident investigation presentation text, create a concise Serious Event Lessons Learned document containing ONLY:
    
    - Brief, clear title/headline
    - Event Summary (brief, readable paragraph)
    - Clearly listed contributing factors
    - Clearly listed lessons learned

    DO NOT INCLUDE sensitive operational details or any unnecessary internal information.

    Here is the presentation text:
    {text}
    
    Format:
    Title:
    Event Summary:
    Contributing Factors:
    - factor 1
    - factor 2
    Lessons Learned:
    - lesson 1
    - lesson 2
    """

    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2,
        max_tokens=700,
    )

    return response.choices[0].message.content.strip()

# Generate Word Document from the summarized content (no template needed)
def create_lessons_learned_doc(content, output_path):
    doc = Document()

    # Parse content into sections
    sections = {}
    current_section = None
    for line in content.split('\n'):
        line = line.strip()
        if line.endswith(':'):
            current_section = line[:-1]
            sections[current_section] = []
        elif current_section:
            sections[current_section].append(line)

    # Title
    title = sections.get("Title", ["Lessons Learned"])[0]
    doc.add_heading(title, level=1)

    # Event Summary
    doc.add_heading("Event Summary", level=2)
    summary = ' '.join(sections.get("Event Summary", []))
    doc.add_paragraph(summary)

    # Contributing Factors
    doc.add_heading("Contributing Factors", level=2)
    for factor in sections.get("Contributing Factors", []):
        doc.add_paragraph(factor, style='List Bullet')

    # Lessons Learned
    doc.add_heading("Lessons Learned", level=2)
    for lesson in sections.get("Lessons Learned", []):
        doc.add_paragraph(lesson, style='List Bullet')

    doc.save(output_path)

# Streamlit UI
st.title('🦺 Serious Event Lessons Learned Generator')

uploaded_file = st.file_uploader("Upload Executive Review PPTX", type=['pptx'])

if uploaded_file:
    input_filepath = f'input/{uploaded_file.name}'
    os.makedirs("input", exist_ok=True)
    with open(input_filepath, 'wb') as f:
        f.write(uploaded_file.getbuffer())

    with st.spinner("Extracting content from presentation..."):
        pptx_text = extract_text_from_pptx(input_filepath)

    with st.spinner("Generating Lessons Learned summary..."):
        generated_content = summarize_and_extract(pptx_text)

    st.success("✅ Generation Complete!")

    st.text_area("📝 Generated Content:", generated_content, height=300)

    output_path = 'generated_lessons_learned.docx'
    create_lessons_learned_doc(generated_content, output_path)

    with open(output_path, "rb") as file:
        st.download_button(
            label="📥 Download Lessons Learned DOCX",
            data=file,
            file_name="Lessons_Learned_Summary.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
